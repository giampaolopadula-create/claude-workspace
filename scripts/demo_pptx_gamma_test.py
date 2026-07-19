"""
Demo: 3 slide di esempio per mostrare il livello grafico raggiungibile
con python-pptx (generazione .pptx nativa, non export da HTML).
Dati e numeri sono ILLUSTRATIVI, solo per mostrare il layout.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
import copy


def add_picture_cover(slide, path, l, t, w, h):
    """Inserisce l'immagine riempiendo il box senza deformarla (crop stile 'object-fit: cover')."""
    img = Image.open(path)
    img_w, img_h = img.size
    box_ratio = w / h
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        crop_frac = 1 - (box_ratio / img_ratio)
        pic = slide.shapes.add_picture(path, l, t, height=h)
        pic.crop_left = crop_frac / 2
        pic.crop_right = crop_frac / 2
        pic.width = w
        pic.left = l
    else:
        crop_frac = 1 - (img_ratio / box_ratio)
        pic = slide.shapes.add_picture(path, l, t, width=w)
        pic.crop_top = crop_frac / 2
        pic.crop_bottom = crop_frac / 2
        pic.height = h
        pic.top = t
    return pic

# ---- Palette luxury (charcoal + gold) ----
CHARCOAL = RGBColor(0x14, 0x14, 0x16)
CHARCOAL_2 = RGBColor(0x1E, 0x1E, 0x21)
GOLD = RGBColor(0xC9, 0xA2, 0x4B)
GOLD_LIGHT = RGBColor(0xE4, 0xCB, 0x94)
IVORY = RGBColor(0xF4, 0xF1, 0xEA)
GREY = RGBColor(0x9A, 0x97, 0x90)

SERIF = "Georgia"
SANS = "Segoe UI Light"
SANS_MED = "Segoe UI Semibold"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

SP = r"C:\Users\GIAMPA~1.PAD\AppData\Local\Temp\claude\c--Users-giampaolo-padula-pig-Desktop-Claude-Workspace\3b0d7895-7b62-4d26-8b16-2c485eb03891\scratchpad"
IMG_COVER = SP + r"\cover_mountain.jpg"
IMG_SIDE = SP + r"\side_mountain.jpg"
IMG_LOBBY = SP + r"\lobby.jpg"


def set_bg(slide, color):
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = color


def add_text(slide, l, t, w, h, text, size, color, font=SANS, bold=False,
             align=PP_ALIGN.LEFT, spacing=None, anchor=MSO_ANCHOR.TOP, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.name = font
    run.font.bold = bold
    run.font.italic = italic
    if spacing:
        rPr = run._r.get_or_add_rPr()
        rPr.set('spc', str(spacing))
    return box


def add_line(slide, l, t, w, h, color, weight=1.0):
    ln = slide.shapes.add_connector(1, l, t, l + w, t + h)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    return ln


def add_rect(slide, l, t, w, h, color, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line_color:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.75)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_overlay(slide, l, t, w, h, color, alpha_pct):
    """Rettangolo semi-trasparente (alpha_pct = opacita' 0-100) per leggibilita' testo su foto."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    alpha_val = str(int(alpha_pct * 1000))
    solid_fill = shp.fill.fore_color._xFill
    srgb = solid_fill.find(qn('a:srgbClr'))
    alpha_el = srgb.makeelement(qn('a:alpha'), {'val': alpha_val})
    srgb.append(alpha_el)
    return shp


# =========================================================
# SLIDE 1 — COVER
# =========================================================
s1 = prs.slides.add_slide(blank)
set_bg(s1, CHARCOAL)

add_picture_cover(s1, IMG_COVER, 0, 0, prs.slide_width, prs.slide_height)
add_overlay(s1, 0, 0, prs.slide_width, prs.slide_height, CHARCOAL, 55)

# sottile cornice dorata
add_line(s1, Inches(0.7), Inches(0.7), Inches(11.933), 0, GOLD, 1.0)
add_line(s1, Inches(0.7), Inches(6.8), Inches(11.933), 0, GOLD, 1.0)

add_text(s1, Inches(0.9), Inches(2.35), Inches(9), Inches(0.5),
          "VALTUR CERVINIA CRISTALLO SKI RESORT", 13, GOLD, SANS_MED,
          spacing=250)

add_text(s1, Inches(0.85), Inches(2.85), Inches(11.5), Inches(1.6),
          "Posizionamento di mercato", 46, IVORY, SERIF)
add_text(s1, Inches(0.85), Inches(3.65), Inches(11.5), Inches(1.6),
          "e leva competitiva 2026/27", 46, IVORY, SERIF, italic=True)

add_line(s1, Inches(0.9), Inches(4.75), Inches(1.6), 0, GOLD, 2.0)

add_text(s1, Inches(0.9), Inches(5.0), Inches(8), Inches(0.5),
          "Executive Summary  —  Direzione Generale", 14, GREY, SANS)

add_text(s1, Inches(0.9), Inches(6.95), Inches(6), Inches(0.4),
          "Luglio 2026", 11, GREY, SANS)
add_text(s1, Inches(9.5), Inches(6.95), Inches(2.9), Inches(0.4),
          "Documento riservato", 11, GREY, SANS, align=PP_ALIGN.RIGHT)

# =========================================================
# SLIDE 2 — KPI DATA SLIDE
# =========================================================
s2 = prs.slides.add_slide(blank)
set_bg(s2, RGBColor(0xFB, 0xFA, 0xF7))

add_text(s2, Inches(0.9), Inches(0.55), Inches(1.5), Inches(0.35),
          "01 — DATI", 12, GOLD, SANS_MED, spacing=200)
add_text(s2, Inches(0.9), Inches(0.9), Inches(7.8), Inches(0.8),
          "Performance stagione in sintesi", 30, CHARCOAL, SERIF)

img_l, img_t, img_w, img_h = Inches(9.6), Inches(0.55), Inches(2.83), Inches(1.5)
add_picture_cover(s2, IMG_LOBBY, img_l, img_t, img_w, img_h)
frame = slide_add_frame = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, img_l, img_t, img_w, img_h)
frame.fill.background()
frame.line.color.rgb = GOLD
frame.line.width = Pt(0.75)
frame.shadow.inherit = False

add_line(s2, Inches(0.9), Inches(1.9), Inches(11.5), 0, RGBColor(0xD8, 0xD3, 0xC8), 1.0)

kpis = [
    ("+7,4%", "RevPAR vs stagione precedente"),
    ("68,2%", "Occupazione media inverno"),
    ("41 gg", "Booking window medio"),
    ("4,7 / 5", "Guest satisfaction score"),
]
col_w = Inches(2.75)
gap = Inches(0.25)
start_l = Inches(0.9)
for i, (num, label) in enumerate(kpis):
    l = start_l + i * (col_w + gap)
    add_text(s2, l, Inches(2.15), col_w, Inches(1.0), num, 40, CHARCOAL, SERIF)
    add_line(s2, l, Inches(3.05), Inches(0.6), 0, GOLD, 2.0)
    add_text(s2, l, Inches(3.25), col_w, Inches(0.8), label, 12.5, GREY, SANS)

# barra comparativa semplice (illustrativa)
add_text(s2, Inches(0.9), Inches(4.35), Inches(8), Inches(0.4),
          "RevPAR per segmento competitivo (indice, illustrativo)", 12.5, CHARCOAL, SANS_MED)

bars = [("Valtur Cervinia", 100, GOLD), ("Comp. set A", 82, RGBColor(0xC9,0xC4,0xB8)),
        ("Comp. set B", 76, RGBColor(0xC9,0xC4,0xB8)), ("Comp. set C", 64, RGBColor(0xC9,0xC4,0xB8))]
bar_top = Inches(4.9)
bar_l = Inches(0.9)
max_w = Inches(9.0)
row_h = Inches(0.55)
for i, (name, val, color) in enumerate(bars):
    t = bar_top + i * row_h
    add_text(s2, bar_l, t, Inches(1.9), row_h, name, 11, CHARCOAL, SANS)
    add_rect(s2, bar_l + Inches(2.0), t + Inches(0.08), Emu(int(max_w * (val/100))), Inches(0.28), color)
    add_text(s2, bar_l + Inches(2.0) + max_w + Inches(0.1), t, Inches(0.6), row_h, str(val), 11, CHARCOAL, SANS_MED)

add_text(s2, Inches(0.9), Inches(7.05), Inches(11), Inches(0.3),
          "Fonte: dati illustrativi a scopo dimostrativo", 9, GREY, SANS, italic=True)

# =========================================================
# SLIDE 3 — TRE COLONNE / CONTENUTO
# =========================================================
s3 = prs.slides.add_slide(blank)
set_bg(s3, CHARCOAL)

# pannello fotografico a destra, a tutta altezza
panel_l = Inches(9.6)
panel_w = prs.slide_width - panel_l
add_picture_cover(s3, IMG_SIDE, panel_l, 0, panel_w, prs.slide_height)
add_overlay(s3, panel_l, 0, panel_w, prs.slide_height, CHARCOAL, 15)
add_line(s3, panel_l, 0, 0, prs.slide_height, GOLD, 1.0)

add_text(s3, Inches(0.9), Inches(0.55), Inches(1.5), Inches(0.35),
          "02 — LEVE", 12, GOLD, SANS_MED, spacing=200)
add_text(s3, Inches(0.9), Inches(0.9), Inches(8.3), Inches(0.8),
          "Tre leve per la stagione 2026/27", 30, IVORY, SERIF)
add_line(s3, Inches(0.9), Inches(1.75), Inches(8.4), 0, RGBColor(0x3A, 0x38, 0x33), 1.0)

cols = [
    ("Direct Booking", "Ridurre la dipendenza dalle OTA rafforzando il canale diretto e il CRM ospiti."),
    ("Experience Premium", "Ampliare l'offerta esperienziale per giustificare l'ADR e allungare la stagione."),
    ("Talent & Servizio", "Investire su formazione e retention dello staff come leva di differenziazione."),
]
row_h3 = Inches(1.55)
start_t3 = Inches(2.15)
for i, (title, body) in enumerate(cols):
    t = start_t3 + i * row_h3
    add_text(s3, Inches(0.9), t, Inches(0.6), Inches(0.5), f"0{i+1}", 20, GOLD, SERIF)
    add_text(s3, Inches(1.75), t, Inches(2.6), Inches(0.4), title, 16, IVORY, SANS_MED)
    add_text(s3, Inches(1.75), t + Inches(0.42), Inches(5.5), Inches(0.9), body, 12, GREY, SANS)

add_line(s3, Inches(0.7), Inches(6.85), Inches(8.6), 0, GOLD, 1.0)
add_text(s3, Inches(0.9), Inches(7.0), Inches(6), Inches(0.4),
          "Valtur Cervinia Cristallo Ski Resort", 10, GREY, SANS)

out = r"c:\Users\giampaolo.padula_pig\Desktop\Claude Workspace\scripts\demo_gamma_test.pptx"
prs.save(out)
print("saved", out)
